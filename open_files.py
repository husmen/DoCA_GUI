""" module for opening various types of files """
import os
import json
import time
import shutil
import subprocess
import textract
import pandas as pd

from docx import Document
#from changeOffice import Change
from pptx import Presentation
from odf.opendocument import load
from odf import text
from pyexcel_xlsx import get_data

class OpenFile():
    def __init__(self, location):
        self.location = location
        print("# opening {}".format(location))
        if (location.endswith("docx") or location.endswith("doc") ):
            self.text, self.paragraphs, self.tables, self.tables2 = self.text_al(location)
        if (location.endswith("pptx") or location.endswith("ppt") or location.endswith("odt") or location.endswith("odp") ):
            self.text, self.paragraphs = self.text_al(location)
            self.tables = None
            self.tables2 = None
        if (location.endswith("xlsx") or location.endswith("xls") or location.endswith("ods")):
            self.tables, self.tables2 = self.text_al(location)

        return

    def text_al(self, dosya_yolu):
        p = []
        t1 = []
        t2 = []
        flag = 0

        if (dosya_yolu.endswith("doc")):
            cwd = os.getcwd()
            #libreoffice --convert-to docx 0020.doc
            dir_path = os.path.dirname(os.path.realpath(dosya_yolu))
            #output_file = dosya_yolu + "x"
            rc = subprocess.call(['libreoffice', '--convert-to', 'docx', '--outdir', dir_path, dosya_yolu])
            output_file = os.path.join(cwd,"tmp/"+os.path.basename(dosya_yolu+"x"))
            for _ in range(5):
                try:    
                    shutil.move(dosya_yolu+"x", output_file)
                except:
                    time.sleep(1)
                else:
                    print("file moved")
                    time.sleep(1)
                    break
            return self.text_al(output_file)
        
        if (dosya_yolu.endswith("xls")):
            cwd = os.getcwd()
            #libreoffice --convert-to docx 0020.doc
            dir_path = os.path.dirname(os.path.realpath(dosya_yolu))
            #output_file = dosya_yolu + "x"
            rc = subprocess.call(['libreoffice', '--convert-to', 'xlsx', '--outdir', dir_path, dosya_yolu])
            output_file = os.path.join(cwd,"tmp/"+os.path.basename(dosya_yolu+"x"))
            for _ in range(5):
                try:    
                    shutil.move(dosya_yolu+"x", output_file)
                except:
                    time.sleep(1)
                else:
                    print("file moved")
                    time.sleep(1)
                    break
            
            return self.text_al(output_file)

        if (dosya_yolu.endswith("ppt")):
            cwd = os.getcwd()
            #libreoffice --convert-to docx 0020.doc
            dir_path = os.path.dirname(os.path.realpath(dosya_yolu))
            #output_file = dosya_yolu + "x"
            rc = subprocess.call(['libreoffice', '--convert-to', 'pptx', '--outdir', dir_path, dosya_yolu])
            output_file = os.path.join(cwd,"tmp/"+os.path.basename(dosya_yolu+"x"))
            for _ in range(5):
                try:    
                    shutil.move(dosya_yolu+"x", output_file)
                except:
                    time.sleep(1)
                else:
                    print("file moved")
                    time.sleep(1)
                    break
            return self.text_al(output_file)

        # docx uzantili
        if (dosya_yolu.endswith("docx")) or (flag == 1):
            doc = Document(dosya_yolu)
            paragraphs = doc.paragraphs
            tables = doc.tables
            fullText = ""
            for paragraph in paragraphs:
                if paragraph != "\n":
                    p.append(paragraph.text)
                    fullText = fullText + paragraph.text + "\n"
            for table in tables:
                t1.append(table)
                #print(table)
                tmp_t = []
                for row in table.rows:
                    tmp_r = []
                    for cell in row.cells:
                        tmp_r.append(cell.text)
                        #print(cell.text)
                    tmp_t.append(tmp_r)
                t2.append(tmp_t)
                #print(tmp_t)
            return fullText, p, t1, t2

        # odt uzantili
        elif (dosya_yolu.endswith("odt")):
            #text1 = textract.process(dosya_yolu)
            doc = load(dosya_yolu)
            paragraphs = []
            txt1 = ""
            for paragraph in doc.getElementsByType(text.P):
                if paragraph != "\n":
                    paragraphs.append(paragraph)
                    txt1 += str(paragraph)

            return txt1, paragraphs
            #print text1
        # pptx uzantili
        elif (dosya_yolu.endswith("pptx")) or (flag == 2):
            paragraphs = []
            prs = Presentation(dosya_yolu)
            text_runs = ""

            for slide in prs.slides:
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    for paragraph in shape.text_frame.paragraphs:
                        tmp = ""
                        for run in paragraph.runs:
                            tmp += run.text
                            text_runs += run.text
                        paragraphs.append(tmp)

            for p in paragraphs:
                if p == "":
                    del p

            return text_runs, paragraphs

        # odp uzantili
        elif (dosya_yolu.endswith("odp")):
            doc = load(dosya_yolu)
            paragraphs = []
            txt1 = ""
            for paragraph in doc.getElementsByType(text.P):
                if paragraph != "\n":
                    paragraphs.append(paragraph)
                    txt1 += str(paragraph)
                    #print (unicode(paragraph))

            return txt1, paragraphs
        # xlsx, xls, ods uzantili
        elif (dosya_yolu.endswith("xlsx")) or (dosya_yolu.endswith("ods")):
            data = get_data(dosya_yolu)
            df_ = []
            #print(data)
            #print(json.dumps(data))
            for sheet in data:
                #print(json.dumps(data[sheet]))
                try:
                    df = pd.DataFrame(data=data[sheet][1:],columns=data[sheet][0])
                except:
                    df = None
                    #df = pd.DataFrame(data[sheet])
                    #df = pd.DataFrame(data=data[sheet][1:][1:],index=data[sheet][1:][0],columns=data[sheet][0][1:])
                df_.append(df)
                #print(df)
            #print(data)
            #print(df_)
            return data, df_